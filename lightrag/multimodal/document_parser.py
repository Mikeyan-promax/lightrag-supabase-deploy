"""
多模态文档解析器

支持从PDF、Word、PPT等文档中提取图像并进行处理
扩展现有文档处理功能，集成图像理解能力
"""

import os
import tempfile
from typing import Dict, List, Optional, Any, Union, Tuple
from pathlib import Path
import logging

logger = logging.getLogger(__name__)

class MultimodalDocumentParser:
    """多模态文档解析器
    
    支持从各种文档格式中提取文本和图像，并进行多模态处理
    """
    
    def __init__(self, image_processor=None):
        """初始化多模态文档解析器
        
        Args:
            image_processor: 图像处理器实例
        """
        self.image_processor = image_processor
        self.supported_formats = {
            'pdf': self._parse_pdf,
            'docx': self._parse_docx,
            'doc': self._parse_doc,
            'pptx': self._parse_pptx,
            'ppt': self._parse_ppt,
            'xlsx': self._parse_xlsx,
            'xls': self._parse_xls
        }
    
    def is_supported_format(self, file_path: Union[str, Path]) -> bool:
        """检查文件格式是否支持
        
        Args:
            file_path: 文件路径
            
        Returns:
            是否支持该格式
        """
        file_ext = Path(file_path).suffix.lower().lstrip('.')
        return file_ext in self.supported_formats
    
    async def parse_document(self, file_path: Union[str, Path]) -> Dict[str, Any]:
        """解析文档，提取文本和图像
        
        Args:
            file_path: 文档文件路径
            
        Returns:
            解析结果，包含文本内容和图像信息
        """
        file_path = Path(file_path)
        if not file_path.exists():
            return {"error": f"文件不存在: {file_path}"}
        
        file_ext = file_path.suffix.lower().lstrip('.')
        if file_ext not in self.supported_formats:
            return {"error": f"不支持的文件格式: {file_ext}"}
        
        try:
            parser_func = self.supported_formats[file_ext]
            result = await parser_func(file_path)
            
            # 处理提取的图像
            if self.image_processor and "images" in result:
                await self._process_extracted_images(result)
            
            return result
            
        except Exception as e:
            logger.error(f"解析文档失败: {e}")
            return {"error": str(e)}
    
    async def _process_extracted_images(self, parse_result: Dict[str, Any]):
        """处理提取的图像
        
        Args:
            parse_result: 解析结果，会被就地修改
        """
        if not self.image_processor or "images" not in parse_result:
            return
        
        processed_images = []
        
        for image_info in parse_result["images"]:
            try:
                if "data" in image_info:
                    # 处理图像数据
                    image_data = image_info["data"]
                    
                    # OCR提取文本
                    ocr_text = await self.image_processor.extract_text(image_data)
                    
                    # 图像内容分析
                    analysis_result = await self.image_processor.analyze_image(
                        image_data, 
                        "请描述这张图片的内容，包括任何可见的文字、图表、图形等信息。"
                    )
                    
                    # 更新图像信息
                    image_info.update({
                        "ocr_text": ocr_text,
                        "analysis": analysis_result,
                        "processed": True
                    })
                
                processed_images.append(image_info)
                
            except Exception as e:
                logger.error(f"处理图像失败: {e}")
                image_info["error"] = str(e)
                processed_images.append(image_info)
        
        parse_result["images"] = processed_images
    
    async def _parse_pdf(self, file_path: Path) -> Dict[str, Any]:
        """解析PDF文档
        
        Args:
            file_path: PDF文件路径
            
        Returns:
            解析结果
        """
        try:
            import fitz  # PyMuPDF
        except ImportError:
            return {"error": "请安装PyMuPDF: pip install PyMuPDF"}
        
        try:
            doc = fitz.open(file_path)
            
            result = {
                "file_path": str(file_path),
                "file_type": "pdf",
                "pages": [],
                "images": [],
                "text_content": "",
                "metadata": {}
            }
            
            # 提取元数据
            result["metadata"] = doc.metadata
            
            # 逐页处理
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                
                # 提取文本
                page_text = page.get_text()
                
                # 提取图像
                page_images = []
                image_list = page.get_images()
                
                for img_index, img in enumerate(image_list):
                    try:
                        # 获取图像数据
                        xref = img[0]
                        pix = fitz.Pixmap(doc, xref)
                        
                        if pix.n - pix.alpha < 4:  # 确保是RGB或灰度图像
                            img_data = pix.tobytes("png")
                            
                            image_info = {
                                "page": page_num + 1,
                                "index": img_index,
                                "data": img_data,
                                "format": "png",
                                "size": (pix.width, pix.height)
                            }
                            
                            page_images.append(image_info)
                            result["images"].append(image_info)
                        
                        pix = None  # 释放内存
                        
                    except Exception as e:
                        logger.warning(f"提取第{page_num+1}页图像{img_index}失败: {e}")
                
                # 保存页面信息
                page_info = {
                    "page_number": page_num + 1,
                    "text": page_text,
                    "images_count": len(page_images)
                }
                
                result["pages"].append(page_info)
                result["text_content"] += page_text + "\n"
            
            doc.close()
            
            result["total_pages"] = len(result["pages"])
            result["total_images"] = len(result["images"])
            
            return result
            
        except Exception as e:
            logger.error(f"解析PDF失败: {e}")
            return {"error": str(e)}
    
    async def _parse_docx(self, file_path: Path) -> Dict[str, Any]:
        """解析DOCX文档
        
        Args:
            file_path: DOCX文件路径
            
        Returns:
            解析结果
        """
        try:
            from docx import Document
            from docx.document import Document as DocumentType
        except ImportError:
            return {"error": "请安装python-docx: pip install python-docx"}
        
        try:
            doc = Document(file_path)
            
            result = {
                "file_path": str(file_path),
                "file_type": "docx",
                "paragraphs": [],
                "images": [],
                "text_content": "",
                "tables": []
            }
            
            # 提取段落文本
            for para in doc.paragraphs:
                if para.text.strip():
                    result["paragraphs"].append(para.text)
                    result["text_content"] += para.text + "\n"
            
            # 提取表格
            for table_index, table in enumerate(doc.tables):
                table_data = []
                for row in table.rows:
                    row_data = [cell.text.strip() for cell in row.cells]
                    table_data.append(row_data)
                
                result["tables"].append({
                    "index": table_index,
                    "data": table_data
                })
            
            # 提取图像
            # 注意：从DOCX中提取图像需要额外的处理
            try:
                import zipfile
                from xml.etree import ElementTree as ET
                
                # DOCX实际上是一个ZIP文件
                with zipfile.ZipFile(file_path, 'r') as zip_file:
                    # 查找媒体文件
                    media_files = [f for f in zip_file.namelist() if f.startswith('word/media/')]
                    
                    for media_file in media_files:
                        try:
                            img_data = zip_file.read(media_file)
                            img_name = os.path.basename(media_file)
                            
                            # 确定图像格式
                            img_format = img_name.split('.')[-1].lower()
                            if img_format in ['jpg', 'jpeg', 'png', 'gif', 'bmp']:
                                image_info = {
                                    "name": img_name,
                                    "data": img_data,
                                    "format": img_format,
                                    "source": "embedded"
                                }
                                
                                result["images"].append(image_info)
                        
                        except Exception as e:
                            logger.warning(f"提取图像{media_file}失败: {e}")
            
            except Exception as e:
                logger.warning(f"提取DOCX图像失败: {e}")
            
            result["total_paragraphs"] = len(result["paragraphs"])
            result["total_images"] = len(result["images"])
            result["total_tables"] = len(result["tables"])
            
            return result
            
        except Exception as e:
            logger.error(f"解析DOCX失败: {e}")
            return {"error": str(e)}
    
    async def _parse_doc(self, file_path: Path) -> Dict[str, Any]:
        """解析DOC文档
        
        Args:
            file_path: DOC文件路径
            
        Returns:
            解析结果
        """
        # DOC格式较复杂，建议转换为DOCX后处理
        return {"error": "DOC格式暂不支持，请转换为DOCX格式"}
    
    async def _parse_pptx(self, file_path: Path) -> Dict[str, Any]:
        """解析PPTX演示文稿
        
        Args:
            file_path: PPTX文件路径
            
        Returns:
            解析结果
        """
        try:
            from pptx import Presentation
        except ImportError:
            return {"error": "请安装python-pptx: pip install python-pptx"}
        
        try:
            prs = Presentation(file_path)
            
            result = {
                "file_path": str(file_path),
                "file_type": "pptx",
                "slides": [],
                "images": [],
                "text_content": ""
            }
            
            # 逐幻灯片处理
            for slide_index, slide in enumerate(prs.slides):
                slide_info = {
                    "slide_number": slide_index + 1,
                    "text": "",
                    "shapes": [],
                    "images_count": 0
                }
                
                # 提取形状和文本
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_info["text"] += shape.text + "\n"
                        slide_info["shapes"].append({
                            "type": "text",
                            "content": shape.text
                        })
                    
                    # 检查是否为图像
                    if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                        try:
                            # 提取图像数据
                            image = shape.image
                            img_data = image.blob
                            
                            image_info = {
                                "slide": slide_index + 1,
                                "data": img_data,
                                "format": image.ext,
                                "size": (shape.width, shape.height)
                            }
                            
                            result["images"].append(image_info)
                            slide_info["images_count"] += 1
                            
                        except Exception as e:
                            logger.warning(f"提取幻灯片{slide_index+1}图像失败: {e}")
                
                result["slides"].append(slide_info)
                result["text_content"] += slide_info["text"]
            
            result["total_slides"] = len(result["slides"])
            result["total_images"] = len(result["images"])
            
            return result
            
        except Exception as e:
            logger.error(f"解析PPTX失败: {e}")
            return {"error": str(e)}
    
    async def _parse_ppt(self, file_path: Path) -> Dict[str, Any]:
        """解析PPT演示文稿
        
        Args:
            file_path: PPT文件路径
            
        Returns:
            解析结果
        """
        # PPT格式较复杂，建议转换为PPTX后处理
        return {"error": "PPT格式暂不支持，请转换为PPTX格式"}
    
    async def _parse_xlsx(self, file_path: Path) -> Dict[str, Any]:
        """解析XLSX电子表格
        
        Args:
            file_path: XLSX文件路径
            
        Returns:
            解析结果
        """
        try:
            import openpyxl
        except ImportError:
            return {"error": "请安装openpyxl: pip install openpyxl"}
        
        try:
            wb = openpyxl.load_workbook(file_path)
            
            result = {
                "file_path": str(file_path),
                "file_type": "xlsx",
                "sheets": [],
                "images": [],
                "text_content": ""
            }
            
            # 逐工作表处理
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                
                sheet_info = {
                    "name": sheet_name,
                    "data": [],
                    "images_count": 0
                }
                
                # 提取单元格数据
                for row in ws.iter_rows(values_only=True):
                    row_data = [str(cell) if cell is not None else "" for cell in row]
                    if any(cell.strip() for cell in row_data):  # 跳过空行
                        sheet_info["data"].append(row_data)
                        result["text_content"] += " ".join(row_data) + "\n"
                
                # 提取图像（如果有）
                if hasattr(ws, '_images'):
                    for img in ws._images:
                        try:
                            img_data = img.ref.getvalue()
                            
                            image_info = {
                                "sheet": sheet_name,
                                "data": img_data,
                                "format": "png",  # 默认格式
                                "anchor": str(img.anchor)
                            }
                            
                            result["images"].append(image_info)
                            sheet_info["images_count"] += 1
                            
                        except Exception as e:
                            logger.warning(f"提取工作表{sheet_name}图像失败: {e}")
                
                result["sheets"].append(sheet_info)
            
            result["total_sheets"] = len(result["sheets"])
            result["total_images"] = len(result["images"])
            
            return result
            
        except Exception as e:
            logger.error(f"解析XLSX失败: {e}")
            return {"error": str(e)}
    
    async def _parse_xls(self, file_path: Path) -> Dict[str, Any]:
        """解析XLS电子表格
        
        Args:
            file_path: XLS文件路径
            
        Returns:
            解析结果
        """
        # XLS格式较复杂，建议转换为XLSX后处理
        return {"error": "XLS格式暂不支持，请转换为XLSX格式"}
    
    async def extract_multimodal_content(self, file_path: Union[str, Path]) -> Dict[str, Any]:
        """提取文档的多模态内容
        
        Args:
            file_path: 文档文件路径
            
        Returns:
            多模态内容，包含文本和图像分析结果
        """
        # 解析文档
        parse_result = await self.parse_document(file_path)
        
        if "error" in parse_result:
            return parse_result
        
        # 整合多模态内容
        multimodal_content = {
            "file_info": {
                "path": parse_result.get("file_path"),
                "type": parse_result.get("file_type"),
                "total_images": parse_result.get("total_images", 0)
            },
            "text_content": parse_result.get("text_content", ""),
            "image_contents": [],
            "combined_content": ""
        }
        
        # 处理图像内容
        if "images" in parse_result:
            for image_info in parse_result["images"]:
                if image_info.get("processed"):
                    image_content = {
                        "ocr_text": image_info.get("ocr_text", ""),
                        "description": image_info.get("analysis", {}).get("content", ""),
                        "metadata": {
                            "page": image_info.get("page"),
                            "slide": image_info.get("slide"),
                            "sheet": image_info.get("sheet"),
                            "size": image_info.get("size")
                        }
                    }
                    multimodal_content["image_contents"].append(image_content)
        
        # 合并所有内容
        combined_parts = []
        
        # 添加文本内容
        if multimodal_content["text_content"].strip():
            combined_parts.append("文档文本内容：\n" + multimodal_content["text_content"])
        
        # 添加图像内容
        for i, img_content in enumerate(multimodal_content["image_contents"]):
            img_part = f"\n图像{i+1}内容：\n"
            
            if img_content["ocr_text"].strip():
                img_part += f"OCR文字：{img_content['ocr_text']}\n"
            
            if img_content["description"].strip():
                img_part += f"图像描述：{img_content['description']}\n"
            
            combined_parts.append(img_part)
        
        multimodal_content["combined_content"] = "\n".join(combined_parts)
        
        return multimodal_content